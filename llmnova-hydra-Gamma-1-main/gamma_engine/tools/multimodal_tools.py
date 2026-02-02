"""Multimodal generation tools for the Gamma Engine.

This module provides tools for generating visual and presentation content,
including images, slides, and diagrams. These tools extend the agent's
capabilities beyond text-based interactions.
"""

import io
import json
import os
import subprocess
from pathlib import Path
from typing import Any, Dict, List, Optional

try:
    from openai import OpenAI
    OPENAI_AVAILABLE = True
except Exception:  # ModuleNotFoundError or similar
    OpenAI = None  # type: ignore
    OPENAI_AVAILABLE = False

import requests  # Needed for ImageGenerationTool, but also useful for general HTTP requests

from .base import Tool

# Optional dependencies - graceful degradation
try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from reportlab.lib.enums import TA_CENTER, TA_LEFT
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import (PageBreak, Paragraph, SimpleDocTemplate,
                                    Spacer)
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

try:
    from PIL import Image
    PILLOW_AVAILABLE = True
except ImportError:
    PILLOW_AVAILABLE = False


class ImageGenerationTool(Tool):
    """Generate images using OpenAI's DALL-E API.
    
    This tool creates images from text descriptions using DALL-E 3,
    OpenAI's advanced image generation model. Generated images are
    saved to the specified output path.
    
    Attributes:
        name: Tool identifier ("generate_image")
        description: Human-readable description
        parameters: JSON Schema for tool parameters
        client: OpenAI API client instance
    
    Examples:
        >>> tool = ImageGenerationTool()
        >>> result = tool.execute(
        ...     prompt="A futuristic city at sunset",
        ...     output_path="./city.png",
        ...     size="1024x1024",
        ...     quality="hd"
        ... )
        >>> print(result)
        Successfully generated image: ./city.png
    """
    
    def __init__(self):
        """Initialize the image generation tool. 
        
        Raises:
            ValueError: If OPENAI_API_KEY environment variable is not set.
        """
        super().__init__(
            name="generate_image",
            description=(
                "Generate an image from a text description using DALL-E 3. "
                "Supports various sizes and quality levels. The generated image "
                "is saved to the specified output path."
            ),
            parameters={
                "type": "object",
                "properties": {
                    "prompt": {
                        "type": "string",
                        "description": "Detailed description of the image to generate"
                    },
                    "size": {
                        "type": "string",
                        "enum": ["1024x1024", "1792x1024", "1024x1792"],
                        "description": "Image dimensions",
                        "default": "1024x1024"
                    },
                    "quality": {
                        "type": "string",
                        "enum": ["standard", "hd"],
                        "description": "Image quality level",
                        "default": "standard"
                    },
                    "output_path": {
                        "type": "string",
                        "description": "Path where the generated image should be saved"
                    }
                },
                "required": ["prompt", "output_path"]
            }
        )
        
        # Initialize OpenAI client only if SDK is available and env var set.
        api_key = os.getenv("OPENAI_API_KEY")
        if not OPENAI_AVAILABLE or not api_key:
            # Do not raise at import/initialization time — keep the class import-safe.
            self.client = None
        else:
            try:
                self.client = OpenAI(api_key=api_key)
            except Exception:
                self.client = None
    
    def execute(
        self,
        prompt: str,
        output_path: str,
        size: str = "1024x1024",
        quality: str = "standard"
    ) -> str:
        """Generate an image from a text prompt. 
        
        Args:
            prompt: Description of the image to generate
            output_path: Where to save the generated image
            size: Image dimensions (1024x1024, 1792x1024, or 1024x1792)
            quality: Quality level (standard or hd)
        
        Returns:
            Success message with path or error message
        
        Raises:
            ValueError: If parameters are invalid
            RuntimeError: If image generation fails
        
        Examples:
            >>> tool = ImageGenerationTool()
            >>> tool.execute(
            ...     prompt="A serene mountain landscape",
            ...     output_path="./mountain.png"
            ... )
            'Successfully generated image: ./mountain.png'
        """
        try:
            # Validate parameters
            if not prompt or len(prompt.strip()) == 0:
                return "Error: Prompt cannot be empty"
            
            if size not in ["1024x1024", "1792x1024", "1024x1792"]:
                return f"Error: Invalid size '{size}'. Must be 1024x1024, 1792x1024, or 1024x1792"
            
            if quality not in ["standard", "hd"]:
                return f"Error: Invalid quality '{quality}'. Must be 'standard' or 'hd'"
            
            # Ensure output directory exists
            output_dir = os.path.dirname(output_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir, exist_ok=True)
            
            # Ensure OpenAI client is available
            if not self.client:
                return (
                    "Error generating image: OpenAI SDK or API key not available. "
                    "Set OPENAI_API_KEY and install openai package to enable image generation."
                )

            # Generate image using DALL-E 3
            response = self.client.images.generate(
                model="dall-e-3",
                prompt=prompt,
                size=size,
                quality=quality,
                n=1,
            )

            # Get the image URL
            image_url = response.data[0].url

            # Download and save the image
            image_response = requests.get(image_url)
            image_response.raise_for_status()

            with open(output_path, "wb") as f:
                f.write(image_response.content)

            return f"Successfully generated image: {output_path}"
            
        except Exception as e:
            return f"Error generating image: {str(e)}"


class SlideGenerationTool(Tool):
    """Generate presentation slides from structured content.
    
    This tool creates professional presentations in PowerPoint (PPTX) or
    PDF format from structured slide data. Supports multiple themes and
    automatic formatting.
    
    Attributes:
        name: Tool identifier ("generate_slides")
        description: Human-readable description
        parameters: JSON Schema for tool parameters
        client: OpenAI API client instance
    
    Examples:
        >>> tool = SlideGenerationTool()
        >>> slides = [
        ...     {"title": "Introduction", "content": ["Point 1", "Point 2"]},
        ...     {"title": "Conclusion", "content": ["Summary"]}
        ... ]
        >>> result = tool.execute(
        ...     title="My Presentation",
        ...     slides=slides,
        ...     output_path="./presentation.pptx",
        ...     format="pptx"
        ... )
    """
    
    def __init__(self):
        """Initialize the slide generation tool."""
        super().__init__(
            name="generate_slides",
            description=(
                "Generate a presentation from structured slide content. "
                "Supports PPTX and PDF output formats with multiple themes. "
                "Each slide can have a title and bullet points or paragraphs."
            ),
            parameters={
                "type": "object",
                "properties": {
                    "title": {
                        "type": "string",
                        "description": "Presentation title (shown on first slide)"
                    },
                    "slides": {
                        "type": "array",
                        "description": "Array of slide objects with title and content",
                        "items": {
                            "type": "object",
                            "properties": {
                                "title": {"type": "string"},
                                "content": {
                                    "type": "array",
                                    "items": {"type": "string"}
                                }
                            },
                            "required": ["title", "content"]
                        }
                    },
                    "output_path": {
                        "type": "string",
                        "description": "Path where the presentation should be saved"
                    },
                    "format": {
                        "type": "string",
                        "enum": ["pptx", "pdf"],
                        "description": "Output format",
                        "default": "pptx"
                    },
                    "theme": {
                        "type": "string",
                        "enum": ["default", "dark", "minimal"],
                        "description": "Presentation theme",
                        "default": "default"
                    }
                },
                "required": ["title", "slides", "output_path"]
            }
        )
    
    def execute(
        self,
        title: str,
        slides: List[Dict[str, Any]],
        output_path: str,
        format: str = "pptx",
        theme: str = "default"
    ) -> str:
        """Generate a presentation from structured content. 
        
        Args:
            title: Presentation title
            slides: List of slide dictionaries with 'title' and 'content' keys
            output_path: Where to save the presentation
            format: Output format ('pptx' or 'pdf')
            theme: Visual theme ('default', 'dark', or 'minimal')
        
        Returns:
            Success message with path or error message
        
        Examples:
            >>> tool = SlideGenerationTool()
            >>> slides = [{"title": "Slide 1", "content": ["Point A", "Point B"]}]
            >>> tool.execute("My Deck", slides, "./deck.pptx")
            'Successfully generated presentation: ./deck.pptx'
        """
        try:
            # Validate parameters
            if not title or len(title.strip()) == 0:
                return "Error: Title cannot be empty"
            
            if not slides or len(slides) == 0:
                return "Error: At least one slide is required"
            
            if format not in ["pptx", "pdf"]:
                return f"Error: Invalid format '{format}'. Must be 'pptx' or 'pdf'"
            
            if theme not in ["default", "dark", "minimal"]:
                return f"Error: Invalid theme '{theme}'. Must be 'default', 'dark', or 'minimal'"
            
            # Ensure output directory exists
            output_dir = os.path.dirname(output_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir, exist_ok=True)
            
            # Generate based on format
            if format == "pptx":
                return self._generate_pptx(title, slides, output_path, theme)
            else:
                return self._generate_pdf(title, slides, output_path, theme)
                
        except Exception as e:
            return f"Error generating slides: {str(e)}"
    
    def _generate_pptx(self, title: str, slides: List[Dict[str, Any]], output_path: str, theme: str) -> str:
        """Generate PowerPoint presentation. 
        
        Args:
            title: Presentation title
            slides: Slide content
            output_path: Output file path
            theme: Visual theme
        
        Returns:
            Success or error message
        """
        if not PPTX_AVAILABLE:
            return "Error: python-pptx not installed. Run: pip install python-pptx"
        
        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Define theme colors
            themes = {
                "default": {"bg": RGBColor(255, 255, 255), "text": RGBColor(0, 0, 0), "accent": RGBColor(0, 112, 192)},
                "dark": {"bg": RGBColor(31, 31, 31), "text": RGBColor(255, 255, 255), "accent": RGBColor(0, 176, 240)},
                "minimal": {"bg": RGBColor(250, 250, 250), "text": RGBColor(51, 51, 51), "accent": RGBColor(100, 100, 100)}
            }
            colors = themes[theme]
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = title
            
            # Apply theme to title slide
            for shape in slide.shapes:
                if hasattr(shape, "fill"):
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = colors["bg"]
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = colors["text"]
            
            # Content slides
            for slide_data in slides:
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)
                
                # Set slide title
                title_shape = slide.shapes.title
                title_shape.text = slide_data["title"]
                title_shape.text_frame.paragraphs[0].font.color.rgb = colors["accent"]
                
                # Set slide background
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = colors["bg"]
                
                # Add content
                content_shape = slide.shapes.placeholders[1]
                text_frame = content_shape.text_frame
                text_frame.clear()
                
                for item in slide_data["content"]:
                    p = text_frame.add_paragraph()
                    p.text = item
                    p.level = 0
                    p.font.size = Pt(18)
                    p.font.color.rgb = colors["text"]
            
            prs.save(output_path)
            return f"Successfully generated presentation: {output_path}"
            
        except Exception as e:
            return f"Error generating PPTX: {str(e)}"
    
    def _generate_pdf(self, title: str, slides: List[Dict[str, Any]], output_path: str, theme: str) -> str:
        """Generate PDF presentation. 
        
        Args:
            title: Presentation title
            slides: Slide content
            output_path: Output file path
            theme: Visual theme
        
        Returns:
            Success or error message
        """
        if not REPORTLAB_AVAILABLE:
            return "Error: reportlab not installed. Run: pip install reportlab"
        
        try:
            doc = SimpleDocTemplate(output_path, pagesize=letter)
            story = []
            styles = getSampleStyleSheet()
            
            # Create custom styles based on theme
            if theme == "dark":
                # Note: ReportLab doesn't easily support dark backgrounds in PDF
                # We'll adjust text colors instead
                title_style = ParagraphStyle(
                    'CustomTitle',
                    parent=styles['Title'],
                    fontSize=36,
                    textColor='#00B0F0',
                    alignment=TA_CENTER,
                    spaceAfter=30
                )
            elif theme == "minimal":
                title_style = ParagraphStyle(
                    'CustomTitle',
                    parent=styles['Title'],
                    fontSize=32,
                    textColor='#333333',
                    alignment=TA_CENTER,
                    spaceAfter=30
                )
            else:
                title_style = ParagraphStyle(
                    'CustomTitle',
                    parent=styles['Title'],
                    fontSize=36,
                    alignment=TA_CENTER,
                    spaceAfter=30
                )
            
            # Title page
            story.append(Spacer(1, 2*inch))
            story.append(Paragraph(title, title_style))
            story.append(PageBreak())
            
            # Content pages
            for slide_data in slides:
                story.append(Paragraph(slide_data["title"], styles['Heading1']))
                story.append(Spacer(1, 0.2*inch))
                
                for item in slide_data["content"]:
                    bullet = f"• {item}"
                    story.append(Paragraph(bullet, styles['Normal']))
                    story.append(Spacer(1, 0.1*inch))
                
                story.append(PageBreak())
            
            doc.build(story)
            return f"Successfully generated presentation: {output_path}"
            
        except Exception as e:
            return f"Error generating PDF: {str(e)}"


class DiagramGenerationTool(Tool):
    """Generate technical diagrams from markup languages.
    
    This tool creates diagrams from Mermaid.js or PlantUML markup,
    supporting flowcharts, sequence diagrams, class diagrams, and more.
    Output formats include PNG and SVG.
    
    Attributes:
        name: Tool identifier ("generate_diagram")
        description: Human-readable description
        parameters: JSON Schema for tool parameters
        client: OpenAI API client instance
    
    Examples:
        >>> tool = DiagramGenerationTool()
        >>> mermaid_code = '''
        ... graph TD
        ...     A[Start] --> B[Process]
        ...     B --> C[End]
        ... '''
        >>> result = tool.execute(
        ...     content=mermaid_code,
        ...     diagram_type="mermaid",
        ...     output_path="./flow.png",
        ...     format="png"
        ... )
    """
    
    def __init__(self):
        """Initialize the diagram generation tool."""
        super().__init__(
            name="generate_diagram",
            description=(
                "Generate technical diagrams from Mermaid.js or PlantUML markup. "
                "Supports flowcharts, sequence diagrams, class diagrams, ER diagrams, "
                "and more. Output can be PNG or SVG format."
            ),
            parameters={
                "type": "object",
                "properties": {
                    "content": {
                        "type": "string",
                        "description": "Diagram definition in Mermaid or PlantUML syntax"
                    },
                    "diagram_type": {
                        "type": "string",
                        "enum": ["mermaid", "plantuml"],
                        "description": "Diagram markup language"
                    },
                    "output_path": {
                        "type": "string",
                        "description": "Path where the diagram should be saved"
                    },
                    "format": {
                        "type": "string",
                        "enum": ["png", "svg"],
                        "description": "Output image format",
                        "default": "png"
                    }
                },
                "required": ["content", "diagram_type", "output_path"]
            }
        )
    
    def execute(
        self,
        content: str,
        diagram_type: str,
        output_path: str,
        format: str = "png"
    ) -> str:
        """Generate a diagram from markup. 
        
        Args:
            content: Diagram definition in Mermaid or PlantUML syntax
            diagram_type: Type of diagram ('mermaid' or 'plantuml')
            output_path: Where to save the diagram
            format: Output format ('png' or 'svg')
        
        Returns:
            Success message with path or error message
        
        Raises:
            ValueError: If parameters are invalid
            RuntimeError: If diagram generation fails
        
        Examples:
            >>> tool = DiagramGenerationTool()
            >>> mermaid = "graph LR\n    A-->B"
            >>> tool.execute(mermaid, "mermaid", "./diagram.png")
            'Successfully generated diagram: ./diagram.png'
        """
        try:
            # Validate parameters
            if not content or len(content.strip()) == 0:
                return "Error: Content cannot be empty"
            
            if diagram_type not in ["mermaid", "plantuml"]:
                return f"Error: Invalid diagram_type '{diagram_type}'. Must be 'mermaid' or 'plantuml'"
            
            if format not in ["png", "svg"]:
                return f"Error: Invalid format '{format}'. Must be 'png' or 'svg'"
            
            # Ensure output directory exists
            output_dir = os.path.dirname(output_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir, exist_ok=True)
            
            # Generate based on diagram type
            if diagram_type == "mermaid":
                return self._generate_mermaid(content, output_path, format)
            else:
                return self._generate_plantuml(content, output_path, format)
                
        except Exception as e:
            return f"Error generating diagram: {str(e)}"
    
    def _generate_mermaid(self, content: str, output_path: str, format: str) -> str:
        """Generate Mermaid diagram. 
        
        Args:
            content: Mermaid markup
            output_path: Output file path
            format: Image format
        
        Returns:
            Success or error message
        """
        try:
            # Create temporary input file
            temp_input = output_path + ".mmd"
            with open(temp_input, 'w', encoding='utf-8') as f:
                f.write(content)
            
            # Try using mmdc (mermaid-cli)
            try:
                cmd = ["mmdc", "-i", temp_input, "-o", output_path]
                if format == "svg":
                    cmd.extend(["-f", "svg"])
                
                result = subprocess.run(
                    cmd,
                    capture_output=True,
                    text=True,
                    timeout=30
                )
                
                # Clean up temp file
                if os.path.exists(temp_input):
                    os.remove(temp_input)
                
                if result.returncode == 0:
                    return f"Successfully generated diagram: {output_path}"
                else:
                    return f"Error: mermaid-cli failed: {result.stderr}"
                    
            except FileNotFoundError:
                # Clean up temp file
                if os.path.exists(temp_input):
                    os.remove(temp_input)
                return (
                    "Error: mermaid-cli (mmdc) not found. "
                    "Install with: npm install -g @mermaid-js/mermaid-cli"
                )
                
        except Exception as e:
            return f"Error generating Mermaid diagram: {str(e)}"
    
    def _generate_plantuml(self, content: str, output_path: str, format: str) -> str:
        """Generate PlantUML diagram. 
        
        Args:
            content: PlantUML markup
            output_path: Output file path
            format: Image format
        
        Returns:
            Success or error message
        """
        try:
            # Create temporary input file
            temp_input = output_path + ".puml"
            with open(temp_input, 'w', encoding='utf-8') as f:
                f.write(content)
            
            # Try using plantuml command
            try:
                format_flag = "-tsvg" if format == "svg" else "-tpng"
                cmd = ["plantuml", format_flag, temp_input, "-o", os.path.dirname(output_path) or "."]
                
                result = subprocess.run(
                    cmd,
                    capture_output=True,
                    text=True,
                    timeout=30
                )
                
                # PlantUML generates with .png or .svg extension
                generated_file = temp_input.replace(".puml", f".{format}")
                
                # Move to desired output path if different
                if generated_file != output_path and os.path.exists(generated_file):
                    import shutil
                    shutil.move(generated_file, output_path)
                
                # Clean up temp file
                if os.path.exists(temp_input):
                    os.remove(temp_input)
                
                if result.returncode == 0 or os.path.exists(output_path):
                    return f"Successfully generated diagram: {output_path}"
                else:
                    return f"Error: plantuml failed: {result.stderr}"
                    
            except FileNotFoundError:
                # Clean up temp file
                if os.path.exists(temp_input):
                    os.remove(temp_input)
                return (
                    "Error: plantuml not found. "
                    "Install Java and download plantuml.jar, or install via package manager"
                )
                
        except Exception as e:
            return f"Error generating PlantUML diagram: {str(e)}"

class AudioGenerationTool(Tool):
    """Generate audio from text (simple TTS shim).

    This is a lightweight, dependency-free fallback tool that converts
    text to a WAV file by using a very small pure-Python approach when
    possible or by writing the text to a file. It's intentionally
    simple so tests can import and instantiate the tool without
    requiring external TTS packages.
    """

    def __init__(self):
        super().__init__(
            name="generate_audio",
            description=(
                "Generate an audio file (WAV) from text. This is a lightweight "
                "shim suitable for tests and demo usage."
            ),
            parameters={
                "type": "object",
                "properties": {
                    "text": {"type": "string", "description": "Text to synthesize"},
                    "output_path": {"type": "string", "description": "Output WAV file path"},
                    "voice": {"type": "string", "description": "Voice identifier (optional)"},
                },
                "required": ["text", "output_path"],
            },
        )

    def execute(self, text: str, output_path: str, voice: Optional[str] = None) -> str:
        """Create a tiny WAV-like file containing the text as bytes.

        This does NOT perform real TTS. It's a safe fallback so the
        tool can be used in unit tests without heavy audio deps.
        """
        try:
            if not text or not text.strip():
                return "Error: text is required"

            # Ensure output directory exists
            output_dir = os.path.dirname(output_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir, exist_ok=True)

            # Write the text to the output path as UTF-8 bytes with a header
            # so the file is meaningful if opened. Keep it small and deterministic.
            with open(output_path, "wb") as f:
                header = b"TTS-SHIM\n"
                f.write(header)
                f.write(text.encode("utf-8"))

            return f"Pseudo-generated audio (text dump) to: {output_path}"
        except Exception as e:
            return f"Error generating audio: {e}"